import os
import requests
import json
import anthropic
import tempfile
import base64
from typing import Optional, List, Dict, Any, BinaryIO, Tuple
import PyPDF2
import docx
from pptx import Presentation
import tiktoken
import openai

# Fonction pour appeler l'API OpenAI
def enhance_description_openai(title: str, initial_description: str, api_key: Optional[str] = None) -> str:
    """
    Améliore la description d'un cours en utilisant l'API OpenAI.
    
    Args:
        title: Le titre du cours
        initial_description: La description initiale fournie par l'utilisateur
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        La description améliorée
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return "Erreur: Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."
    
    # Préparer la requête pour l'API OpenAI
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    
    # Construire le prompt pour l'API
    prompt = f"""
    En tant qu'expert en pédagogie, améliore cette description de cours:
    
    Titre du cours: {title}
    
    Description initiale: {initial_description}
    
    Veuillez fournir une description améliorée qui:
    1. Est professionnelle et engageante
    2. Inclut les objectifs d'apprentissage clairs
    3. Mentionne les bénéfices pour l'apprenant
    4. Est structurée en paragraphes concis
    5. Fait environ 150-200 mots
    """
    
    # Données pour l'API
    data = {
        "model": "gpt-4o",
        "messages": [
            {"role": "system", "content": "Vous êtes un expert en pédagogie et en création de contenu éducatif."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7,
        "max_tokens": 5000
    }
    
    try:
        # Appel à l'API
        response = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers=headers,
            data=json.dumps(data)
        )
        
        # Vérifier si la requête a réussi
        response.raise_for_status()
        
        # Extraire la réponse
        result = response.json()
        enhanced_description = result["choices"][0]["message"]["content"].strip()
        
        return enhanced_description
    
    except Exception as e:
        return f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"

# Fonction pour appeler l'API Anthropic
def enhance_description_anthropic(title: str, initial_description: str, api_key: Optional[str] = None) -> str:
    """
    Améliore la description d'un cours en utilisant l'API Anthropic Claude.
    
    Args:
        title: Le titre du cours
        initial_description: La description initiale fournie par l'utilisateur
        api_key: Clé API Anthropic (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        La description améliorée
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("ANTHROPIC_API_KEY")
    
    if not api_key:
        return "Erreur: Clé API Anthropic non trouvée. Veuillez configurer la variable d'environnement ANTHROPIC_API_KEY."
    
    try:
        # Créer le client Anthropic avec la bibliothèque officielle
        client = anthropic.Anthropic(api_key=api_key)
        
        # Construire le message utilisateur
        user_prompt = f"""En tant qu'expert en pédagogie, améliore cette description de cours:
        
Titre du cours: {title}

Description initiale: {initial_description}

Veuillez fournir une description améliorée qui:
1. Est professionnelle et engageante
2. Inclut les objectifs d'apprentissage clairs
3. Mentionne les bénéfices pour l'apprenant
4. Est structurée en paragraphes concis
5. Fait environ 150-200 mots"""
        
        # Appel à l'API avec la bibliothèque officielle
        message = client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=4000,
            system="Vous êtes un expert en pédagogie et en création de contenu éducatif.",
            messages=[
                {"role": "user", "content": user_prompt}
            ]
        )
        
        # Extraire la réponse
        enhanced_description = message.content[0].text
        
        return enhanced_description
    
    except Exception as e:
        return f"Erreur lors de l'appel à l'API Anthropic: {str(e)}"

# Fonction principale qui choisit l'API à utiliser
def enhance_course_description(title: str, initial_description: str, api_provider: str = "openai") -> str:
    """
    Améliore la description d'un cours en utilisant l'API spécifiée.
    
    Args:
        title: Le titre du cours
        initial_description: La description initiale fournie par l'utilisateur
        api_provider: Le fournisseur d'API à utiliser ('openai' ou 'anthropic')
    
    Returns:
        La description améliorée
    """
    if not title or not initial_description:
        return "Veuillez fournir un titre et une description initiale."
    
    if api_provider.lower() == "openai":
        return enhance_description_openai(title, initial_description)
    elif api_provider.lower() == "anthropic":
        return enhance_description_anthropic(title, initial_description)
    else:
        return f"Fournisseur d'API non pris en charge: {api_provider}. Utilisez 'openai' ou 'anthropic'."

# Nouvelles fonctions pour générer du contenu avec OpenAI

def generate_learning_objectives(course_title: str, course_description: str, api_key: Optional[str] = None) -> List[str]:
    """
    Génère des objectifs d'apprentissage pour un cours en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        course_description: La description du cours
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Liste d'objectifs d'apprentissage
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return ["Erreur: Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."]
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Construire le prompt
    prompt = f"""
    En tant qu'expert en pédagogie, génère 5 objectifs d'apprentissage clairs et mesurables pour ce cours:
    
    Titre du cours: {course_title}
    
    Description du cours: {course_description}
    
    Les objectifs d'apprentissage doivent:
    1. Commencer par des verbes d'action (ex: identifier, analyser, créer)
    2. Être spécifiques et mesurables
    3. Être alignés avec le contenu du cours
    4. Couvrir différents niveaux de la taxonomie de Bloom (connaissance, compréhension, application, analyse, synthèse, évaluation)
    5. Être formulés du point de vue de l'apprenant
    
    Format de sortie: Une liste numérotée de 5 objectifs d'apprentissage.
    """
    
    try:
        # Appel à l'API
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "Vous êtes un expert en pédagogie et en création de contenu éducatif."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1000
        )
        
        # Extraire la réponse
        objectives_text = response.choices[0].message.content.strip()
        
        # Traiter la réponse pour extraire les objectifs
        objectives = []
        for line in objectives_text.split('\n'):
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith('- ')):
                # Supprimer les numéros ou puces au début
                cleaned_line = line.lstrip('0123456789.- ')
                objectives.append(cleaned_line)
        
        # Si aucun objectif n'a été extrait, retourner le texte brut
        if not objectives:
            objectives = [objectives_text]
        
        return objectives
    
    except Exception as e:
        return [f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"]

def generate_prerequisites(course_title: str, course_description: str, api_key: Optional[str] = None) -> List[str]:
    """
    Génère des prérequis pour un cours en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        course_description: La description du cours
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Liste de prérequis
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return ["Erreur: Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."]
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Construire le prompt
    prompt = f"""
    En tant qu'expert en pédagogie, génère une liste de prérequis pour ce cours:
    
    Titre du cours: {course_title}
    
    Description du cours: {course_description}
    
    Les prérequis doivent:
    1. Être clairs et spécifiques
    2. Inclure les connaissances et compétences nécessaires
    3. Mentionner les outils ou logiciels requis, le cas échéant
    4. Indiquer le niveau d'expérience préalable recommandé
    5. Être réalistes et pertinents pour le cours
    
    Format de sortie: Une liste de 4-6 prérequis.
    """
    
    try:
        # Appel à l'API
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "Vous êtes un expert en pédagogie et en création de contenu éducatif."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1000
        )
        
        # Extraire la réponse
        prerequisites_text = response.choices[0].message.content.strip()
        
        # Traiter la réponse pour extraire les prérequis
        prerequisites = []
        for line in prerequisites_text.split('\n'):
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith('- ')):
                # Supprimer les numéros ou puces au début
                cleaned_line = line.lstrip('0123456789.- ')
                prerequisites.append(cleaned_line)
        
        # Si aucun prérequis n'a été extrait, retourner le texte brut
        if not prerequisites:
            prerequisites = [prerequisites_text]
        
        return prerequisites
    
    except Exception as e:
        return [f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"]

def generate_learning_methods(course_title: str, course_description: str, api_key: Optional[str] = None) -> List[str]:
    """
    Génère des méthodes d'apprentissage pour un cours en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        course_description: La description du cours
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Liste de méthodes d'apprentissage
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return ["Erreur: Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."]
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Construire le prompt
    prompt = f"""
    En tant qu'expert en pédagogie, génère 3-5 méthodes d'apprentissage efficaces pour ce cours:
    
    Titre du cours: {course_title}
    
    Description du cours: {course_description}
    
    Les méthodes d'apprentissage doivent:
    1. Être variées et adaptées au contenu du cours
    2. Inclure des approches pédagogiques modernes et efficaces
    3. Être clairement nommées (ex: "Vidéos interactives", "Ateliers pratiques", "Études de cas")
    4. Être pertinentes pour le sujet du cours
    
    Format de sortie: Une liste de 3-5 méthodes d'apprentissage.
    """
    
    try:
        # Appel à l'API
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "Vous êtes un expert en pédagogie et en création de contenu éducatif."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1000
        )
        
        # Extraire la réponse
        methods_text = response.choices[0].message.content.strip()
        
        # Traiter la réponse pour extraire les méthodes
        methods = []
        for line in methods_text.split('\n'):
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith('- ')):
                # Supprimer les numéros ou puces au début
                cleaned_line = line.lstrip('0123456789.- ')
                methods.append(cleaned_line)
        
        # Si aucune méthode n'a été extraite, retourner le texte brut
        if not methods:
            methods = [methods_text]
        
        return methods
    
    except Exception as e:
        return [f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"]

# Fonctions pour traiter différents types de documents

def extract_text_from_pdf(file: BinaryIO) -> str:
    """
    Extrait le texte d'un fichier PDF.
    
    Args:
        file: Le fichier PDF en mode binaire
    
    Returns:
        Le texte extrait du PDF
    """
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + "\n\n"
        return text
    except Exception as e:
        return f"Erreur lors de l'extraction du texte du PDF: {str(e)}"

def extract_text_from_docx(file: BinaryIO) -> str:
    """
    Extrait le texte d'un fichier Word (DOCX).
    
    Args:
        file: Le fichier DOCX en mode binaire
    
    Returns:
        Le texte extrait du document Word
    """
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
            temp_file.write(file.read())
            temp_file_path = temp_file.name
        
        doc = docx.Document(temp_file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        # Supprimer le fichier temporaire
        os.unlink(temp_file_path)
        
        return text
    except Exception as e:
        return f"Erreur lors de l'extraction du texte du document Word: {str(e)}"

def extract_text_from_pptx(file: BinaryIO) -> str:
    """
    Extrait le texte d'un fichier PowerPoint (PPTX).
    
    Args:
        file: Le fichier PPTX en mode binaire
    
    Returns:
        Le texte extrait de la présentation PowerPoint
    """
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(file.read())
            temp_file_path = temp_file.name
        
        prs = Presentation(temp_file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        # Supprimer le fichier temporaire
        os.unlink(temp_file_path)
        
        return text
    except Exception as e:
        return f"Erreur lors de l'extraction du texte de la présentation PowerPoint: {str(e)}"

def extract_text_from_txt(file: BinaryIO) -> str:
    """
    Extrait le texte d'un fichier texte (TXT).
    
    Args:
        file: Le fichier TXT en mode binaire
    
    Returns:
        Le texte extrait du fichier texte
    """
    try:
        content = file.read()
        # Essayer de décoder avec différents encodages
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        
        # Si aucun encodage ne fonctionne, utiliser utf-8 avec des erreurs ignorées
        return content.decode('utf-8', errors='ignore')
    except Exception as e:
        return f"Erreur lors de l'extraction du texte du fichier texte: {str(e)}"

def split_text_into_chunks(text: str, max_tokens: int = 8000) -> List[str]:
    """
    Divise un texte en morceaux plus petits pour respecter les limites de tokens.
    
    Args:
        text: Le texte à diviser
        max_tokens: Le nombre maximum de tokens par morceau
    
    Returns:
        Une liste de morceaux de texte
    """
    try:
        # Initialiser l'encodeur de tokens pour GPT-4
        enc = tiktoken.encoding_for_model("gpt-4")
        
        # Diviser le texte en paragraphes
        paragraphs = text.split("\n")
        
        chunks = []
        current_chunk = ""
        current_token_count = 0
        
        for paragraph in paragraphs:
            # Compter les tokens dans ce paragraphe
            paragraph_tokens = len(enc.encode(paragraph))
            
            # Si ajouter ce paragraphe dépasse la limite, commencer un nouveau morceau
            if current_token_count + paragraph_tokens > max_tokens:
                if current_chunk:  # Éviter d'ajouter des morceaux vides
                    chunks.append(current_chunk)
                current_chunk = paragraph
                current_token_count = paragraph_tokens
            else:
                # Sinon, ajouter le paragraphe au morceau actuel
                if current_chunk:
                    current_chunk += "\n" + paragraph
                else:
                    current_chunk = paragraph
                current_token_count += paragraph_tokens
        
        # Ajouter le dernier morceau s'il n'est pas vide
        if current_chunk:
            chunks.append(current_chunk)
        
        return chunks
    except Exception as e:
        print(f"Erreur lors de la division du texte: {str(e)}")
        # En cas d'erreur, diviser simplement par caractères
        chunks = []
        chunk_size = 4000  # Approximation grossière
        for i in range(0, len(text), chunk_size):
            chunks.append(text[i:i + chunk_size])
        return chunks

def create_embeddings(text_chunks: List[str], api_key: Optional[str] = None) -> List[List[float]]:
    """
    Crée des embeddings pour une liste de morceaux de texte en utilisant l'API OpenAI.
    
    Args:
        text_chunks: Liste de morceaux de texte
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Liste d'embeddings (vecteurs)
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        raise ValueError("Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY.")
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    embeddings = []
    
    for chunk in text_chunks:
        try:
            response = client.embeddings.create(
                model="text-embedding-3-small",
                input=chunk
            )
            embedding = response.data[0].embedding
            embeddings.append(embedding)
        except Exception as e:
            print(f"Erreur lors de la création de l'embedding: {str(e)}")
            # En cas d'erreur, ajouter un embedding vide
            embeddings.append([0.0] * 1536)  # Dimension standard pour les embeddings OpenAI
    
    return embeddings

def process_document(file: BinaryIO, filename: str, api_key: Optional[str] = None) -> Tuple[str, List[List[float]]]:
    """
    Traite un document, extrait son texte et crée des embeddings.
    
    Args:
        file: Le fichier en mode binaire
        filename: Le nom du fichier
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Un tuple contenant le texte extrait et les embeddings
    """
    # Déterminer le type de fichier
    file_extension = os.path.splitext(filename)[1].lower()
    
    # Extraire le texte en fonction du type de fichier
    if file_extension == '.pdf':
        text = extract_text_from_pdf(file)
    elif file_extension == '.docx':
        text = extract_text_from_docx(file)
    elif file_extension == '.pptx':
        text = extract_text_from_pptx(file)
    elif file_extension == '.txt':
        text = extract_text_from_txt(file)
    else:
        return f"Type de fichier non pris en charge: {file_extension}", []
    
    # Diviser le texte en morceaux
    text_chunks = split_text_into_chunks(text)
    
    # Créer des embeddings
    embeddings = create_embeddings(text_chunks, api_key)
    
    return text, embeddings

def generate_course_structure(
    course_title: str, 
    course_description: str, 
    duration: str, 
    difficulty: str, 
    num_modules: int,
    document_text: str = "",
    document_embeddings: List[List[float]] = [],
    api_key: Optional[str] = None
) -> Dict[str, Any]:
    """
    Génère une structure hiérarchique de modules et chapitres pour un cours en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        course_description: La description du cours
        duration: La durée du cours
        difficulty: Le niveau de difficulté du cours
        num_modules: Le nombre de modules souhaité
        document_text: Le texte des documents uploadés (optionnel)
        document_embeddings: Les embeddings des documents uploadés (optionnel)
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Un dictionnaire contenant la structure du cours avec modules et chapitres
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return {"error": "Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."}
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Construire le contexte pour l'IA
    context = f"""
    Titre du cours: {course_title}
    Description du cours: {course_description}
    Durée: {duration}
    Niveau de difficulté: {difficulty}
    Nombre de modules souhaité: {num_modules}
    """
    
    # Ajouter le contenu des documents si disponible
    if document_text:
        context += f"\nContenu des documents de référence: {document_text[:2000]}..."
    
    # Construire le prompt pour l'API
    prompt = f"""
    En tant qu'expert en pédagogie et en conception de cours, crée une structure hiérarchique complète pour un cours avec le contexte suivant:
    
    {context}
    
    Génère une structure de cours avec:
    1. {num_modules} modules numérotés et titrés
    2. Pour chaque module, 3 à 5 chapitres numérotés et titrés
    3. Pour chaque chapitre, une brève description du contenu (2-3 phrases)
    4. Pour chaque chapitre, 3 à 5 points clés qui seront abordés
    
    La structure doit être progressive, cohérente et adaptée au niveau de difficulté indiqué.
    Assure-toi que les modules et chapitres s'enchaînent logiquement et couvrent l'ensemble du sujet.
    
    Retourne la structure sous forme d'un objet JSON structuré comme suit:
    {{
        "modules": [
            {{
                "module_number": 1,
                "module_title": "Titre du module 1",
                "chapters": [
                    {{
                        "chapter_number": 1.1,
                        "chapter_title": "Titre du chapitre 1.1",
                        "description": "Description du chapitre 1.1",
                        "key_points": ["Point clé 1", "Point clé 2", "Point clé 3"]
                    }},
                    ...
                ]
            }},
            ...
        ]
    }}
    """
    
    try:
        # Appel à l'API OpenAI
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "Tu es un expert en pédagogie et en conception de cours. Tu dois créer une structure de cours détaillée et cohérente."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000,
            response_format={"type": "json_object"}
        )
        
        # Extraire et parser la réponse JSON
        course_structure_text = response.choices[0].message.content
        course_structure = json.loads(course_structure_text)
        
        return course_structure
    
    except Exception as e:
        return {"error": f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"}

def generate_chapter_content(
    course_title: str,
    course_description: str,
    module_title: str,
    chapter_title: str,
    chapter_description: str,
    key_points: List[str],
    document_text: str = "",
    document_embeddings: List[List[float]] = [],
    api_key: Optional[str] = None
) -> Dict[str, Any]:
    """
    Génère le contenu détaillé d'un chapitre de cours en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        course_description: La description du cours
        module_title: Le titre du module
        chapter_title: Le titre du chapitre
        chapter_description: La description du chapitre
        key_points: Les points clés du chapitre
        document_text: Le texte des documents uploadés (optionnel)
        document_embeddings: Les embeddings des documents uploadés (optionnel)
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Un dictionnaire contenant le contenu détaillé du chapitre
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return {"error": "Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."}
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Construire le contexte pour l'IA
    context = f"""
    Titre du cours: {course_title}
    Description du cours: {course_description}
    Titre du module: {module_title}
    Titre du chapitre: {chapter_title}
    Description du chapitre: {chapter_description}
    Points clés à aborder:
    {', '.join(key_points)}
    """
    
    # Ajouter le contenu des documents si disponible
    if document_text:
        context += f"\nContenu des documents de référence: {document_text[:2000]}..."
    
    # Construire le prompt pour l'API
    prompt = f"""
    En tant qu'expert en pédagogie et en conception de cours, crée un contenu détaillé pour un chapitre de cours avec le contexte suivant:
    
    {context}
    
    Génère un contenu de chapitre complet qui inclut:
    1. Une introduction engageante qui présente le sujet du chapitre
    2. Un développement structuré qui couvre tous les points clés mentionnés
    3. Des exemples concrets et des illustrations pour faciliter la compréhension
    4. Des explications claires et pédagogiques adaptées au niveau de difficulté du cours
    5. Une conclusion qui résume les points importants et fait le lien avec le reste du cours
    6. Des questions de réflexion ou exercices pratiques pour renforcer l'apprentissage
    
    Le contenu doit être informatif, engageant et pédagogiquement solide.
    
    Retourne le contenu sous forme d'un objet JSON structuré comme suit:
    {{
        "introduction": "Texte d'introduction",
        "sections": [
            {{
                "title": "Titre de la section 1",
                "content": "Contenu détaillé de la section 1",
                "examples": ["Exemple 1", "Exemple 2"]
            }},
            ...
        ],
        "conclusion": "Texte de conclusion",
        "exercises": [
            {{
                "question": "Question 1",
                "answer": "Réponse ou indice pour la question 1"
            }},
            ...
        ]
    }}
    """
    
    try:
        # Appel à l'API OpenAI
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "Tu es un expert en pédagogie et en conception de cours. Tu dois créer un contenu de chapitre détaillé, informatif et pédagogiquement solide."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000,
            response_format={"type": "json_object"}
        )
        
        # Extraire et parser la réponse JSON
        chapter_content_text = response.choices[0].message.content
        chapter_content = json.loads(chapter_content_text)
        
        return chapter_content
    
    except Exception as e:
        return {"error": f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"}

def generate_quiz(
    course_title: str,
    module_data: Dict[str, Any],
    num_questions: int = 10,
    difficulty_level: str = "Moyen",
    question_types: List[str] = ["Choix multiple", "Vrai/Faux", "Questions directes"],
    api_key: Optional[str] = None
) -> Dict[str, Any]:
    """
    Génère un quiz basé sur le contenu d'un module en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        module_data: Les données du module (titre, chapitres, etc.)
        num_questions: Le nombre de questions à générer
        difficulty_level: Le niveau de difficulté du quiz
        question_types: Les types de questions à inclure
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Un dictionnaire contenant le quiz généré
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return {"error": "Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."}
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Extraire les informations du module
    module_title = module_data["module_title"]
    module_number = module_data["module_number"]
    
    # Construire un résumé du contenu du module
    module_content = f"Module {module_number}: {module_title}\n\n"
    
    # Ajouter le contenu des chapitres
    for chapter in module_data["chapters"]:
        module_content += f"Chapitre {chapter['chapter_number']}: {chapter['chapter_title']}\n"
        module_content += f"Description: {chapter['description']}\n"
        module_content += "Points clés:\n"
        for point in chapter['key_points']:
            module_content += f"- {point}\n"
        module_content += "\n"
    
    # Construire le prompt pour l'API
    prompt = f"""
    En tant qu'expert en pédagogie, crée un quiz pour évaluer les connaissances sur le module suivant:
    
    Cours: {course_title}
    {module_content}
    
    Génère un quiz avec les caractéristiques suivantes:
    - Nombre de questions: {num_questions}
    - Niveau de difficulté: {difficulty_level}
    - Types de questions à inclure: {', '.join(question_types)}
    
    Pour chaque question, inclus:
    1. L'énoncé de la question
    2. Le type de question (parmi ceux spécifiés)
    3. Les options de réponse (pour les questions à choix multiple et vrai/faux)
    4. La réponse correcte
    5. Une explication de la réponse
    
    Le quiz doit couvrir équitablement l'ensemble du contenu du module et être adapté au niveau de difficulté demandé.
    
    Retourne le quiz sous forme d'un objet JSON structuré comme suit:
    {{
        "module_title": "Titre du module",
        "module_number": X,
        "quiz_title": "Titre du quiz",
        "difficulty_level": "Niveau de difficulté",
        "questions": [
            {{
                "question_number": 1,
                "question_text": "Énoncé de la question",
                "question_type": "Type de question",
                "options": ["Option A", "Option B", "Option C", "Option D"],  // Pour les questions à choix multiple
                "correct_answer": "Réponse correcte",
                "explanation": "Explication de la réponse"
            }},
            ...
        ]
    }}
    """
    
    try:
        # Appel à l'API OpenAI
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "Tu es un expert en pédagogie et en création de quiz. Tu dois créer un quiz pertinent et adapté au contenu fourni."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000,
            response_format={"type": "json_object"}
        )
        
        # Extraire et parser la réponse JSON
        quiz_text = response.choices[0].message.content
        quiz = json.loads(quiz_text)
        
        return quiz
    
    except Exception as e:
        return {"error": f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"}

def generate_podcast_script(
    course_title: str,
    course_description: str,
    course_structure: Dict[str, Any],
    podcast_format: str = "Interview",
    podcast_duration: str = "15-20 minutes",
    target_audience: str = "Étudiants",
    api_key: Optional[str] = None
) -> Dict[str, Any]:
    """
    Génère un script de podcast basé sur le contenu du cours en utilisant l'API OpenAI.
    
    Args:
        course_title: Le titre du cours
        course_description: La description du cours
        course_structure: La structure du cours (modules et chapitres)
        podcast_format: Le format du podcast (Interview, Monologue, Discussion, etc.)
        podcast_duration: La durée cible du podcast
        target_audience: Le public cible du podcast
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Un dictionnaire contenant le script du podcast
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return {"error": "Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."}
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    # Construire un résumé du contenu du cours
    course_content = f"Titre du cours: {course_title}\n"
    course_content += f"Description: {course_description}\n\n"
    
    # Ajouter un résumé de la structure du cours
    course_content += "Structure du cours:\n"
    
    if "modules" in course_structure:
        for module in course_structure["modules"]:
            course_content += f"Module {module['module_number']}: {module['module_title']}\n"
            
            # Ajouter un résumé des chapitres du module
            for chapter in module["chapters"]:
                course_content += f"  - Chapitre {chapter['chapter_number']}: {chapter['chapter_title']}\n"
                course_content += f"    Description: {chapter['description']}\n"
                course_content += "    Points clés:\n"
                for point in chapter['key_points']:
                    course_content += f"      * {point}\n"
    
    # Construire le prompt pour l'API
    prompt = f"""
    En tant qu'expert en création de contenu audio pédagogique, crée un script de podcast sur le cours suivant:
    
    {course_content}
    
    Caractéristiques du podcast:
    - Format: {podcast_format}
    - Durée cible: {podcast_duration}
    - Public cible: {target_audience}
    
    Le script doit inclure:
    1. Une introduction engageante qui présente le sujet du cours
    2. Une présentation claire des principaux concepts et idées du cours
    3. Des discussions sur les points clés de chaque module
    4. Des exemples concrets et des anecdotes pour illustrer les concepts
    5. Une conclusion qui résume les points importants et encourage l'auditeur à en apprendre davantage
    
    Si le format est "Interview" ou "Discussion", inclure des dialogues entre un hôte et un ou plusieurs invités experts.
    Si le format est "Monologue", structurer le contenu comme une narration continue par un seul présentateur.
    
    Le script doit être conversationnel, engageant et adapté à un format audio.
    
    Retourne le script sous forme d'un objet JSON structuré comme suit:
    {{
        "podcast_title": "Titre du podcast",
        "format": "Format du podcast",
        "duration": "Durée estimée",
        "target_audience": "Public cible",
        "participants": ["Nom du participant 1", "Nom du participant 2", ...],
        "script_sections": [
            {{
                "section_title": "Introduction",
                "content": "Texte du script pour cette section"
            }},
            {{
                "section_title": "Développement - Module X",
                "content": "Texte du script pour cette section"
            }},
            ...
            {{
                "section_title": "Conclusion",
                "content": "Texte du script pour cette section"
            }}
        ]
    }}
    """
    
    try:
        # Appel à l'API OpenAI
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "Tu es un expert en création de contenu audio pédagogique. Tu dois créer un script de podcast informatif, engageant et adapté au format audio."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000,
            response_format={"type": "json_object"}
        )
        
        # Extraire et parser la réponse JSON
        podcast_script_text = response.choices[0].message.content
        podcast_script = json.loads(podcast_script_text)
        
        return podcast_script
    
    except Exception as e:
        return {"error": f"Erreur lors de l'appel à l'API OpenAI: {str(e)}"}

def generate_podcast_audio(
    script_text: str,
    voice: str = "alloy",
    api_key: Optional[str] = None
) -> Dict[str, Any]:
    """
    Convertit un script de podcast en audio en utilisant l'API OpenAI TTS (Text-to-Speech).
    
    Args:
        script_text: Le texte du script à convertir en audio
        voice: La voix à utiliser pour la synthèse vocale (alloy, echo, fable, onyx, nova, shimmer)
        api_key: Clé API OpenAI (optionnelle, sinon utilise la variable d'environnement)
    
    Returns:
        Un dictionnaire contenant l'URL du fichier audio et d'autres métadonnées
    """
    # Utiliser la clé API fournie ou celle de l'environnement
    api_key = api_key or os.environ.get("OPENAI_API_KEY")
    
    if not api_key:
        return {"error": "Clé API OpenAI non trouvée. Veuillez configurer la variable d'environnement OPENAI_API_KEY."}
    
    # Initialiser le client OpenAI
    client = openai.OpenAI(api_key=api_key)
    
    try:
        # Limiter la longueur du texte si nécessaire (l'API TTS a des limites)
        max_length = 4096  # Limite approximative pour l'API TTS
        if len(script_text) > max_length:
            script_text = script_text[:max_length]
        
        # Appel à l'API OpenAI TTS
        response = client.audio.speech.create(
            model="tts-1",
            voice=voice,
            input=script_text
        )
        
        # Obtenir directement les données binaires de la réponse
        audio_data = response.content
        
        # Encoder en base64 pour pouvoir l'utiliser dans Streamlit
        audio_base64 = base64.b64encode(audio_data).decode("utf-8")
        
        return {
            "success": True,
            "audio_base64": audio_base64,
            "duration": len(audio_data) / 16000,  # Estimation approximative de la durée en secondes
            "format": "mp3"
        }
    
    except Exception as e:
        return {"error": f"Erreur lors de la génération de l'audio: {str(e)}"} 